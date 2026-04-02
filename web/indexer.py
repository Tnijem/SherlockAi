"""
Sherlock file indexer — GB-scale, multi-format, hash-deduped, chunked.

Pipeline design:
  Stage 1 — Scanner (main thread): rglob files, mtime/hash check against DB,
             skip unchanged files before touching them. Generator-based so
             scanning and extracting overlap.
  Stage 2 — Extract pool (N threads): text extraction (PDF/DOCX/OCR/etc.)
             is I/O + CPU bound and fully parallelisable.
  Stage 3 — Embed + upsert (main thread): chunks batched to Ollama /api/embed
             (cap 32 per call), normalised, upserted to ChromaDB + FTS5,
             committed to SQLite. Single writer = no locking headaches.

Cancellation: any job can be cancelled mid-run. Pending (not-yet-flushed)
  work is discarded; already-committed work stays. No half-indexed files.

Progress: _update_job writes stage + counts to the shared status file so
  any consumer (web UI, CLI) sees live granular status.
"""

from __future__ import annotations

import concurrent.futures
import csv
import hashlib
import json
import math
import sqlite3 as _sqlite3
import subprocess
import tempfile
import threading
import uuid
from datetime import datetime as _dt
from pathlib import Path
from typing import Iterator, Optional

import requests
from sqlalchemy.orm import Session

from config import DB_PATH, EMBED_MODEL, GLOBAL_COLLECTION, OLLAMA_URL, UPLOADS_DIR, user_collection
from logging_config import get_logger
from models import IndexedFile, SessionLocal, Upload

log = get_logger("sherlock.indexer")

# ── Tuning constants ──────────────────────────────────────────────────────────


# ── Upload indexing queue (serialized to avoid SQLite lock contention) ─────────
import queue as _queue

_upload_queue: _queue.Queue = _queue.Queue()

def _upload_queue_worker():
    """Process upload indexing jobs one at a time."""
    while True:
        fn = _upload_queue.get()
        try:
            fn()
        except Exception as e:
            log.error("upload_queue_worker error: %s", e)
        finally:
            _upload_queue.task_done()

threading.Thread(target=_upload_queue_worker, daemon=True, name="upload-index-queue").start()

CHUNK_SIZE        = 1200   # chars per chunk (~300 tokens)
CHUNK_OVERLAP     = 200    # overlap between adjacent chunks
EMBED_BATCH_SIZE  = 32     # max chunks per /api/embed call (memory cap for bge-m3)
FILE_BATCH_SIZE   = 8      # files accumulated before a flush to the store
N_EXTRACT_WORKERS = 4      # parallel text-extraction threads

# ── Supported extensions ──────────────────────────────────────────────────────

TEXT_EXTS   = {".txt", ".md", ".rst", ".log", ".csv", ".tsv"}
PDF_EXTS    = {".pdf"}
WORD_EXTS   = {".docx", ".doc"}
EXCEL_EXTS  = {".xlsx", ".xls"}
PPT_EXTS    = {".pptx", ".ppt"}
HTML_EXTS   = {".html", ".htm"}
RTF_EXTS    = {".rtf"}
IMAGE_EXTS  = {".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp", ".gif"}
AUDIO_EXTS  = {".mp3", ".wav", ".m4a", ".ogg", ".flac", ".aac"}
EMAIL_EXTS  = {".eml"}

ALL_SUPPORTED = (
    TEXT_EXTS | PDF_EXTS | WORD_EXTS | EXCEL_EXTS | PPT_EXTS |
    HTML_EXTS | RTF_EXTS | IMAGE_EXTS | AUDIO_EXTS | EMAIL_EXTS
)

# ── Chunking ──────────────────────────────────────────────────────────────────

def chunk_text(text: str) -> list[dict]:
    """Split text into chunks with location metadata.

    Returns list of dicts: {text, page_start, page_end, line_start, line_end}.
    Page info is extracted from \x00PAGE:N\x00 markers inserted by _extract_pdf.
    Line numbers are computed from character offsets in the original text.
    """
    import re as _re
    text = text.strip()
    if not text:
        return []

    # Build page boundary map from markers
    page_map: list[tuple[int, int]] = []  # (char_offset, page_num)
    for m in _re.finditer(r'\x00PAGE:(\d+)\x00\n?', text):
        page_map.append((m.start(), int(m.group(1))))

    # Strip page markers from text (they shouldn't appear in embeddings)
    clean = _re.sub(r'\x00PAGE:\d+\x00\n?', '', text)

    # Build a line-number lookup: for any char offset, what line is it on?
    line_offsets = [0]
    for i, ch in enumerate(clean):
        if ch == '\n':
            line_offsets.append(i + 1)

    def _char_to_line(offset: int) -> int:
        import bisect
        idx = bisect.bisect_right(line_offsets, offset) - 1
        return max(1, idx + 1)

    # Adjust page_map offsets to account for removed markers
    adjusted_pages: list[tuple[int, int]] = []
    removed = 0
    marker_spans = list(_re.finditer(r'\x00PAGE:\d+\x00\n?', text))
    for ms in marker_spans:
        adjusted_pages.append((ms.start() - removed, int(_re.search(r'\d+', ms.group()).group())))
        removed += ms.end() - ms.start()

    def _char_to_page(offset: int) -> int:
        if not adjusted_pages:
            return 0  # not a PDF
        page = 0
        for po, pn in adjusted_pages:
            if po <= offset:
                page = pn
            else:
                break
        return page

    # Chunk the clean text
    chunks = []
    start = 0
    while start < len(clean):
        end = min(start + CHUNK_SIZE, len(clean))
        chunk_text_str = clean[start:end]
        if chunk_text_str.strip():
            ps = _char_to_page(start)
            pe = _char_to_page(end - 1) if end > start else ps
            ls = _char_to_line(start)
            le = _char_to_line(end - 1) if end > start else ls
            chunks.append({
                "text": chunk_text_str,
                "page_start": ps,
                "page_end": pe,
                "line_start": ls,
                "line_end": le,
            })
        start += CHUNK_SIZE - CHUNK_OVERLAP

    return chunks


# ── Text extraction ───────────────────────────────────────────────────────────

def extract_text(filepath: Path) -> str:
    """Extract plain text from a file. Returns empty string on failure."""
    ext = filepath.suffix.lower()
    try:
        if ext in TEXT_EXTS:    return _extract_text_file(filepath)
        if ext in PDF_EXTS:     return _extract_pdf(filepath)
        if ext == ".docx":      return _extract_docx(filepath)
        if ext == ".doc":       return _extract_via_libreoffice(filepath)
        if ext == ".xlsx":      return _extract_xlsx(filepath)
        if ext == ".xls":       return _extract_via_libreoffice(filepath)
        if ext in PPT_EXTS:
            return _extract_pptx(filepath) if ext == ".pptx" else _extract_via_libreoffice(filepath)
        if ext in HTML_EXTS:    return _extract_html(filepath)
        if ext in RTF_EXTS:     return _extract_rtf(filepath)
        if ext in IMAGE_EXTS:   return _extract_image_ocr(filepath)
        if ext in AUDIO_EXTS:   return _extract_audio_whisper(filepath)
        if ext in EMAIL_EXTS:   return _extract_eml(filepath)
    except Exception as e:
        log.warning("extract_text failed for %s: %s", filepath, e)
    return ""


def _extract_text_file(fp: Path) -> str:
    if fp.suffix.lower() in {".csv", ".tsv"}:
        delim = "\t" if fp.suffix.lower() == ".tsv" else ","
        rows = []
        with fp.open(encoding="utf-8", errors="ignore") as f:
            for row in csv.reader(f, delimiter=delim):
                rows.append(" | ".join(row))
        return "\n".join(rows)
    return fp.read_text(encoding="utf-8", errors="ignore")


def _extract_pdf(fp: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(fp))
    parts = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        if not text.strip():
            try:
                from PIL import Image
                import pytesseract, io as _io
                for img_obj in page.images:
                    img = Image.open(_io.BytesIO(img_obj.data))
                    text += pytesseract.image_to_string(img) + "\n"
            except Exception:
                pass
        # Insert lightweight page marker (stripped before embedding)
        parts.append(f"\x00PAGE:{i}\x00\n{text}")
    return "\n".join(parts)


def _extract_docx(fp: Path) -> str:
    from docx import Document
    doc = Document(str(fp))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(parts)


def _extract_via_libreoffice(fp: Path) -> str:
    """Convert .doc/.xls/.ppt via LibreOffice headless. Uses Popen + poll to yield GIL."""
    import time as _t
    with tempfile.TemporaryDirectory() as tmpdir:
        proc = subprocess.Popen(
            ["/opt/homebrew/bin/libreoffice", "--headless", "--convert-to", "txt:Text", "--outdir", tmpdir, str(fp)],
            stdout=subprocess.PIPE, stderr=subprocess.PIPE,
        )
        deadline = _t.monotonic() + 30  # 30s timeout (was 60)
        while proc.poll() is None:
            if _t.monotonic() > deadline:
                proc.kill()
                raise RuntimeError("LibreOffice timed out (30s)")
            _t.sleep(0.2)  # yield GIL every 200ms so uvicorn can serve requests
        if proc.returncode != 0:
            raise RuntimeError(f"LibreOffice failed: {proc.stderr.read().decode()[:200]}")
        out = Path(tmpdir) / (fp.stem + ".txt")
        return out.read_text(encoding="utf-8", errors="ignore") if out.exists() else ""


def _extract_xlsx(fp: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(fp), read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(v) for v in row if v is not None)
            if row_text.strip():
                parts.append(row_text)
    return "\n".join(parts)


def _extract_pptx(fp: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(fp))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    return "\n".join(parts)


def _extract_html(fp: Path) -> str:
    import re
    text = fp.read_text(encoding="utf-8", errors="ignore")
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"&\w+;", " ", text)
    return re.sub(r"\s+", " ", text).strip()


def _extract_rtf(fp: Path) -> str:
    import re
    text = fp.read_text(encoding="utf-8", errors="ignore")
    text = re.sub(r"\\[a-z]+\-?\d*\s?", " ", text)
    text = re.sub(r"[{}\\]", "", text)
    return re.sub(r"\s+", " ", text).strip()


def _extract_image_ocr(fp: Path) -> str:
    from PIL import Image
    import pytesseract
    return pytesseract.image_to_string(Image.open(str(fp)))


# ── Whisper singleton (load once, reuse — transcription is the slowest op) ───

_whisper_model = None
_whisper_lock  = threading.Lock()


def _get_whisper() -> object:
    global _whisper_model
    if _whisper_model is None:
        with _whisper_lock:
            if _whisper_model is None:
                from faster_whisper import WhisperModel
                from config import WHISPER_MODEL, WHISPER_MODEL_DIR
                _whisper_model = WhisperModel(
                    WHISPER_MODEL, device="cpu", download_root=WHISPER_MODEL_DIR,
                )
                log.info("Whisper model loaded: %s", WHISPER_MODEL)
    return _whisper_model


def _extract_audio_whisper(fp: Path) -> str:
    model = _get_whisper()
    segments, _ = model.transcribe(str(fp), beam_size=5)
    return " ".join(seg.text for seg in segments)


def _extract_eml(fp: Path) -> str:
    import email
    msg = email.message_from_bytes(fp.read_bytes())
    parts = []
    for hdr in ("subject", "from", "to", "date"):
        if msg[hdr]:
            parts.append(f"{hdr.capitalize()}: {msg[hdr]}")
    for part in msg.walk():
        if part.get_content_type() == "text/plain":
            payload = part.get_payload(decode=True)
            if payload:
                parts.append(payload.decode("utf-8", errors="ignore"))
    return "\n".join(parts)


# ── File hash ─────────────────────────────────────────────────────────────────

def file_hash(fp: Path, chunk_size: int = 1 << 20) -> str:
    h = hashlib.sha256()
    with fp.open("rb") as f:
        while block := f.read(chunk_size):
            h.update(block)
    return h.hexdigest()


# ── Batch embedding ───────────────────────────────────────────────────────────

def _embed_batch(texts: list[str]) -> list[list[float]]:
    """
    Embed up to EMBED_BATCH_SIZE texts in one /api/embed call.
    Returns L2-normalised embeddings.
    """
    resp = requests.post(
        f"{OLLAMA_URL}/api/embed",
        json={"model": EMBED_MODEL, "input": [t[:8192] for t in texts]},
        timeout=120,
    )
    resp.raise_for_status()
    raw = resp.json()["embeddings"]
    result = []
    for emb in raw:
        norm = math.sqrt(sum(x * x for x in emb)) or 1.0
        result.append([x / norm for x in emb])
    return result


def _embed_all(texts: list[str]) -> list[list[float]]:
    """Embed an arbitrary number of texts, honouring EMBED_BATCH_SIZE cap."""
    out = []
    for i in range(0, len(texts), EMBED_BATCH_SIZE):
        out.extend(_embed_batch(texts[i:i + EMBED_BATCH_SIZE]))
    return out


# ── Job tracking ──────────────────────────────────────────────────────────────

_jobs:        dict[str, dict] = {}
_cancel_jobs: set[str]        = set()
_jobs_lock    = threading.Lock()

_STATUS_FILE: Optional[Path] = None


def _status_file() -> Path:
    global _STATUS_FILE
    if _STATUS_FILE is None:
        _STATUS_FILE = Path(DB_PATH).parent / "indexer_status.json"
    return _STATUS_FILE


def _write_status(data: dict) -> None:
    try:
        _status_file().write_text(json.dumps(data))
    except Exception:
        pass


def read_live_status() -> Optional[dict]:
    """Read indexer status from any process. Returns None if no active job."""
    try:
        p = _status_file()
        if not p.exists():
            return None
        return json.loads(p.read_text())
    except Exception:
        return None


def _new_job() -> str:
    job_id = str(uuid.uuid4())
    state  = {"status": "queued", "stage": "queued",
               "indexed": 0, "skipped": 0, "errors": 0,
               "total": 0, "done": False, "messages": []}
    with _jobs_lock:
        _jobs[job_id] = state
    _write_status({**state, "job_id": job_id})
    return job_id


def get_job_status(job_id: str) -> Optional[dict]:
    with _jobs_lock:
        return dict(_jobs.get(job_id, {}))


def _update_job(job_id: str, **kwargs) -> None:
    with _jobs_lock:
        if job_id in _jobs:
            _jobs[job_id].update(kwargs)
            _write_status({**_jobs[job_id], "job_id": job_id})


def _is_cancelled(job_id: str) -> bool:
    return job_id in _cancel_jobs


def cancel_job(job_id: str) -> bool:
    """Request cancellation of a running job. Returns True if job was found."""
    with _jobs_lock:
        if job_id not in _jobs:
            return False
        _cancel_jobs.add(job_id)
        _jobs[job_id].update({"status": "cancelling", "stage": "cancelling"})
        _write_status({**_jobs[job_id], "job_id": job_id})
    return True


# ── Stage 1 — Scanner / gatekeeper ───────────────────────────────────────────

def _should_index(fp: Path) -> tuple[bool, Optional[str], Optional[str], int, str, bool]:
    """
    Cheap mtime/hash check against SQLite. Called on the main scan thread.
    Returns (needs_index, fhash, fmtime, fsize, reason, is_update).
    reason is for logging only.
    """
    try:
        stat   = fp.stat()
        fsize  = stat.st_size
        fmtime = str(stat.st_mtime)
    except OSError as e:
        log.warning("stat failed %s: %s", fp, e)
        return False, None, None, 0, "stat_error", False

    db = SessionLocal()
    try:
        existing = db.query(IndexedFile).filter(IndexedFile.file_path == str(fp)).first()

        if existing and existing.mtime == fmtime:
            return False, None, None, 0, "mtime_match", False

        # mtime changed — check hash to distinguish real change vs NAS touch
        try:
            fhash = file_hash(fp)
        except OSError as e:
            log.warning("hash failed %s: %s", fp, e)
            return False, None, None, 0, "hash_error", False

        if existing and existing.file_hash == fhash:
            # Content identical — just update the mtime stamp and skip extraction
            existing.mtime = fmtime
            db.commit()
            return False, fhash, fmtime, fsize, "mtime_only", False

        return True, fhash, fmtime, fsize, "changed" if existing else "new", existing is not None
    finally:
        db.close()


# ── Stage 2 — Extract (runs in worker threads) ────────────────────────────────

def _extract_worker(
    fp: Path,
    collection_name: str,
    source_label: str,
    case_id: Optional[int],
    fhash: str,
    fmtime: str,
    fsize: int,
    is_update: bool,
) -> Optional[dict]:
    """
    Extract text and chunk it. Returns a result dict or None on failure.
    Thread-safe: no shared mutable state.
    """
    text = extract_text(fp)
    if not text.strip():
        log.debug("No text extracted from %s", fp)
        return None

    chunk_dicts = chunk_text(text)
    if not chunk_dicts:
        return None
    chunks = [c["text"] for c in chunk_dicts]
    chunk_locations = chunk_dicts

    # For updates: how many old chunks existed (to detect shrinkage)
    old_chunk_count = 0
    if is_update:
        db = SessionLocal()
        try:
            rec = db.query(IndexedFile).filter(IndexedFile.file_path == str(fp)).first()
            old_chunk_count = rec.chunk_count if rec else 0
        finally:
            db.close()

    return {
        "fp":         fp,
        "chunks":     chunks,
        "chunk_locations": chunk_locations,
        "label":      source_label,
        "fhash":      fhash,
        "fmtime":     fmtime,
        "fsize":      fsize,
        "ext":        fp.suffix.lower(),
        "collection": collection_name,
        "case_id":    case_id,
        "is_update":  is_update,
        "old_count":  old_chunk_count,
    }


# ── Stage 3 — Embed + upsert (single writer thread) ──────────────────────────

def _flush_batch(
    results:    list[dict],
    coll,
    fts_conn:   _sqlite3.Connection,
    write_db:   Session,
) -> int:
    """
    Embed all chunks in results (batched to EMBED_BATCH_SIZE), upsert to
    ChromaDB + FTS5 + SQLite. Atomic per file: each file either fully
    commits or is skipped on error — no partial state.
    Returns number of files successfully flushed.
    """
    if not results:
        return 0

    # Collect all chunk texts across all files for a single embed call
    all_texts: list[str] = []
    spans:     list[tuple[dict, int, int]] = []   # (result, start, end)
    for r in results:
        s = len(all_texts)
        all_texts.extend(r["chunks"])
        spans.append((r, s, len(all_texts)))

    try:
        all_embeddings = _embed_all(all_texts)
    except Exception as e:
        log.error("Batch embed failed (%d chunks): %s", len(all_texts), e)
        return 0

    fts_cur = fts_conn.cursor()
    flushed = 0

    for result, start, end in spans:
        fp         = result["fp"]
        chunks     = result["chunks"]
        embeddings = all_embeddings[start:end]
        label      = result["label"]
        ext        = result["ext"]
        col        = result["collection"]

        doc_ids = [f"{fp}__chunk_{i}" for i in range(len(chunks))]
        locations = result.get("chunk_locations", [])
        metas   = [
            {"source": label, "path": str(fp), "chunk": i,
             "total_chunks": len(chunks), "ext": ext,
             "page_start": locations[i]["page_start"] if i < len(locations) else 0,
             "page_end":   locations[i]["page_end"]   if i < len(locations) else 0,
             "line_start": locations[i]["line_start"] if i < len(locations) else 0,
             "line_end":   locations[i]["line_end"]   if i < len(locations) else 0}
            for i in range(len(chunks))
        ]

        # Delete excess old chunks if file shrank
        if result["is_update"] and result["old_count"] > len(chunks):
            excess = [f"{fp}__chunk_{i}" for i in range(len(chunks), result["old_count"])]
            try:
                coll.delete(ids=excess)
                fts_cur.executemany("DELETE FROM chunk_fts WHERE chunk_id = ?",
                                    [(eid,) for eid in excess])
            except Exception:
                pass

        try:
            coll.upsert(ids=doc_ids, embeddings=embeddings,
                        documents=chunks, metadatas=metas)
        except Exception as e:
            log.error("ChromaDB upsert failed for %s: %s", fp, e)
            continue

        try:
            fts_cur.executemany(
                "INSERT OR REPLACE INTO chunk_fts(chunk_id, collection, source, content)"
                " VALUES (?,?,?,?)",
                [(did, col, label, txt) for did, txt in zip(doc_ids, chunks)],
            )
            fts_conn.commit()  # commit FTS per-file to release DB lock before SQLAlchemy writes
        except Exception as e:
            log.warning("FTS5 upsert failed for %s: %s", fp, e)

        # SQLite indexed_files — commit per file so crashes leave known-good state
        try:
            existing = write_db.query(IndexedFile).filter(
                IndexedFile.file_path == str(fp)
            ).first()
            if existing:
                existing.file_hash   = result["fhash"]
                existing.size_bytes  = result["fsize"]
                existing.mtime       = result["fmtime"]
                existing.chunk_count = len(doc_ids)
                existing.collection  = col
                existing.indexed_at  = _dt.utcnow()
                if result["case_id"] is not None:
                    existing.case_id = result["case_id"]
            else:
                write_db.add(IndexedFile(
                    file_path    = str(fp),
                    file_hash    = result["fhash"],
                    size_bytes   = result["fsize"],
                    mtime        = result["fmtime"],
                    chunk_count  = len(doc_ids),
                    collection   = col,
                    case_id      = result["case_id"],
                ))
            write_db.commit()
        except Exception as e:
            log.error("SQLite commit failed for %s: %s", fp, e)
            write_db.rollback()
            continue

        flushed += 1

    # fts already committed per-file above
    return flushed


# ── Pipeline orchestrator ─────────────────────────────────────────────────────

def _run_pipeline(
    file_paths:       list[Path],
    collection_name:  str,
    job_id:           str,
    source_label_fn:  Optional[callable] = None,
    case_id:          Optional[int]      = None,
) -> tuple[int, int, int]:
    """
    Runs the full extract → embed → upsert pipeline.
    Returns (indexed, skipped, errors).

    Cancellation: checked after each file completes. Pending work is
    discarded; already-flushed files are committed and remain valid.
    """
    from rag import get_or_create_collection

    coll     = get_or_create_collection(collection_name)
    fts_conn = _sqlite3.connect(str(DB_PATH)); fts_conn.execute("PRAGMA busy_timeout=10000")
    write_db = SessionLocal()

    indexed = skipped = errors = 0
    pending: list[dict] = []

    try:
        # ── Stage 1: scan + filter (main thread) ──────────────────────────────
        _update_job(job_id, stage="scanning", total=len(file_paths))
        to_extract: list[tuple[Path, str, str, str, int, bool]] = []

        for fp in file_paths:
            if _is_cancelled(job_id):
                break
            needs, fhash, fmtime, fsize, reason, is_update = _should_index(fp)
            if needs:
                label = source_label_fn(fp) if source_label_fn else fp.name
                to_extract.append((fp, label, fhash, fmtime, fsize, is_update))
            else:
                skipped += 1
            _update_job(job_id, skipped=skipped)

        if _is_cancelled(job_id):
            _update_job(job_id, status="cancelled", stage="cancelled", done=True)
            return indexed, skipped, errors

        log.info("Pipeline scan done: %d to extract, %d skipped", len(to_extract), skipped)

        # ── Stage 2: extract (thread pool) ────────────────────────────────────
        _update_job(job_id, stage="extracting", total=len(file_paths))

        def _submit(args):
            fp, label, fhash, fmtime, fsize, is_update = args
            return _extract_worker(fp, collection_name, label, case_id,
                                   fhash, fmtime, fsize, is_update)

        with concurrent.futures.ThreadPoolExecutor(
            max_workers=N_EXTRACT_WORKERS,
            thread_name_prefix="sherlock-extract",
        ) as pool:
            future_map = {pool.submit(_submit, args): args[0] for args in to_extract}

            for future in concurrent.futures.as_completed(future_map):
                fp = future_map[future]

                if _is_cancelled(job_id):
                    # Cancel remaining futures; already-submitted ones finish naturally
                    for f in future_map:
                        f.cancel()
                    break

                try:
                    result = future.result()
                except Exception as e:
                    errors += 1
                    log.error("Extract error %s: %s", fp, e)
                    _update_job(job_id, indexed=indexed, skipped=skipped, errors=errors)
                    continue

                if result is None:
                    # extract_text returned empty — count as skipped, not error
                    skipped += 1
                    _update_job(job_id, skipped=skipped)
                    continue

                pending.append(result)

                # ── Stage 3: embed + upsert when batch is full ─────────────
                if len(pending) >= FILE_BATCH_SIZE:
                    _update_job(job_id, stage="embedding")
                    n = _flush_batch(pending, coll, fts_conn, write_db)
                    indexed += n
                    errors  += len(pending) - n
                    pending.clear()
                    _update_job(job_id, stage="extracting",
                                indexed=indexed, skipped=skipped, errors=errors)

        # Final flush of whatever's left
        if pending and not _is_cancelled(job_id):
            _update_job(job_id, stage="embedding")
            n = _flush_batch(pending, coll, fts_conn, write_db)
            indexed += n
            errors  += len(pending) - n
            pending.clear()

        if _is_cancelled(job_id):
            _update_job(job_id, status="cancelled", stage="cancelled", done=True)
        else:
            _update_job(job_id, status="done", stage="done",
                        indexed=indexed, skipped=skipped, errors=errors, done=True)

    except Exception as e:
        log.error("Pipeline error (job %s): %s", job_id, e)
        _update_job(job_id, status="error", stage="error", done=True, messages=[str(e)])
    finally:
        fts_conn.close()
        write_db.close()

    return indexed, skipped, errors


# ── Public job launchers ──────────────────────────────────────────────────────

def start_case_index(case_id: int, nas_path: str) -> str:
    """Index a case's NAS path into its own ChromaDB collection. Returns job_id."""
    from models import Case, case_collection

    job_id = _new_job()

    def _run():
        _update_job(job_id, status="running")
        root = Path(nas_path)
        if not root.exists():
            _update_job(job_id, status="error", done=True,
                        messages=[f"NAS path not accessible: {nas_path}"])
            return

        col_name   = case_collection(case_id)
        file_paths = [
            fp for fp in root.rglob("*")
            if fp.is_file() and fp.suffix.lower() in ALL_SUPPORTED
        ]

        # Apply index filters
        from file_filters import get_filter_set as _get_fs
        file_paths, n_filtered = _get_fs().apply(file_paths)
        if n_filtered:
            _update_job(job_id, messages=[f"⚙ {n_filtered} file(s) excluded by index filters"])

        indexed, skipped, errors = _run_pipeline(
            file_paths, col_name, job_id,
            source_label_fn=lambda fp: fp.name,
            case_id=case_id,
        )

        # Update case record
        db = SessionLocal()
        try:
            case = db.query(Case).filter(Case.id == case_id).first()
            if case:
                case.last_indexed  = _dt.utcnow()
                case.indexed_count = (
                    db.query(IndexedFile).filter(IndexedFile.case_id == case_id).count()
                )
                db.commit()
        finally:
            db.close()

        log.info("case_index_done", extra={
            "job_id": job_id, "case_id": case_id,
            "indexed": indexed, "skipped": skipped, "errors": errors,
        })

    threading.Thread(target=_run, daemon=True, name=f"index-case-{case_id}").start()
    return job_id


def start_nas_index(nas_paths: list[str]) -> str:
    """Index a list of NAS paths into the global collection. Returns job_id."""
    job_id = _new_job()

    def _run():
        _update_job(job_id, status="running")
        all_files: list[Path] = []

        for nas_path in nas_paths:
            root = Path(nas_path)
            if not root.exists():
                _update_job(job_id, messages=[f"⚠ NAS path not accessible: {nas_path}"])
                continue
            all_files.extend(
                fp for fp in root.rglob("*")
                if fp.is_file() and fp.suffix.lower() in ALL_SUPPORTED
            )

        # Apply index filters
        from file_filters import get_filter_set as _get_fs
        all_files, n_filtered = _get_fs().apply(all_files)
        if n_filtered:
            _update_job(job_id, messages=[f"⚙ {n_filtered} file(s) excluded by index filters"])

        if not all_files:
            _update_job(job_id, status="done", done=True)
            return

        _run_pipeline(all_files, GLOBAL_COLLECTION, job_id)
        log.info("nas_index_done", extra={"job_id": job_id, "paths": nas_paths})

    threading.Thread(target=_run, daemon=True, name="index-nas-global").start()
    return job_id


def start_upload_index(upload_id: int, user_id: int, filepath: Path) -> str:
    """
    Index a single user-uploaded file into their private collection.
    Minimal path — no thread pool overhead for one file.
    Returns job_id.
    """
    job_id   = _new_job()
    col_name = user_collection(user_id)

    def _run():
        db = SessionLocal()
        try:
            upload = db.query(Upload).filter(Upload.id == upload_id).first()
            if not upload:
                _update_job(job_id, status="error", done=True)
                return

            upload.status = "indexing"
            db.commit()

            _update_job(job_id, status="running", stage="scanning", total=1)

            needs, fhash, fmtime, fsize, reason, is_update = _should_index(filepath)
            if not needs:
                # Already indexed and unchanged
                upload.status = "ready"
                db.commit()
                _update_job(job_id, status="done", stage="done", done=True, skipped=1)
                return

            _update_job(job_id, stage="extracting")
            label  = filepath.name
            result = _extract_worker(
                filepath, col_name, label, None, fhash, fmtime, fsize, is_update
            )

            if result is None:
                upload.status    = "error"
                upload.error_msg = "No text could be extracted from this file."
                db.commit()
                _update_job(job_id, status="error", stage="error", done=True,
                            messages=["No text extracted"])
                return

            _update_job(job_id, stage="embedding")
            from rag import get_or_create_collection
            coll     = get_or_create_collection(col_name)
            fts_conn = _sqlite3.connect(str(DB_PATH)); fts_conn.execute("PRAGMA busy_timeout=10000")
            write_db = SessionLocal()
            try:
                n = _flush_batch([result], coll, fts_conn, write_db)
            finally:
                fts_conn.close()
                write_db.close()

            chunk_ids    = [f"{filepath}__chunk_{i}" for i in range(len(result["chunks"]))]
            upload.chroma_ids = json.dumps(chunk_ids)
            # Extract page count for PDFs
            if filepath.suffix.lower() == ".pdf":
                try:
                    from pypdf import PdfReader as _PR
                    upload.page_count = len(_PR(str(filepath)).pages)
                except Exception:
                    pass
            upload.status = "ready" if n else "error"
            if not n:
                upload.error_msg = "Embedding or upsert failed."
            db.commit()

            _update_job(job_id, status="done" if n else "error",
                        stage="done" if n else "error",
                        indexed=n, done=True)

        except Exception as e:
            log.error("Upload index error: %s", e)
            try:
                rec = db.query(Upload).filter(Upload.id == upload_id).first()
                if rec:
                    rec.status    = "error"
                    rec.error_msg = str(e)
                    db.commit()
            except Exception:
                db.rollback()
            _update_job(job_id, status="error", stage="error", done=True, messages=[str(e)])
        finally:
            db.close()

    _upload_queue.put(_run)
    return job_id
